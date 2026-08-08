"""Binary exporters: report document -> PDF / Word / PowerPoint bytes.

Each exporter is import-guarded so a missing optional library degrades to a
clear error instead of breaking the API. Layouts follow the shared document
structure (cover, table of contents, numbered sections, tables, callouts) so
all formats read like the same professionally produced report.
"""
from __future__ import annotations

import io
from datetime import datetime
from typing import Any

ACCENT = "#2563EB"
INK = "#16233B"
MUTED = "#5B6B84"
LINE = "#E3E8F0"
BG = "#F7F9FC"


class ExportUnavailable(RuntimeError):
    """Raised when the library backing a format is not installed."""


def available_formats() -> dict[str, bool]:
    out = {"markdown": True, "html": True}
    for fmt, module in (("pdf", "reportlab"), ("docx", "docx"), ("pptx", "pptx")):
        try:
            __import__(module)
            out[fmt] = True
        except ImportError:
            out[fmt] = False
    return out


def _date(doc: dict[str, Any]) -> str:
    try:
        return datetime.fromisoformat(doc.get("generated_at", "")).strftime("%B %d, %Y")
    except ValueError:
        return doc.get("generated_at", "")


# --- PDF (reportlab / platypus) ----------------------------------------------
def render_pdf(doc: dict[str, Any]) -> bytes:
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import mm
        from reportlab.platypus import (
            PageBreak,
            Paragraph,
            SimpleDocTemplate,
            Spacer,
            Table,
            TableStyle,
        )
    except ImportError as exc:  # pragma: no cover
        raise ExportUnavailable("PDF export requires the 'reportlab' package.") from exc

    accent = colors.HexColor(ACCENT)
    ink = colors.HexColor(INK)
    muted = colors.HexColor(MUTED)
    line = colors.HexColor(LINE)
    bg = colors.HexColor(BG)

    styles = getSampleStyleSheet()
    body = ParagraphStyle("Body", parent=styles["Normal"], fontSize=9.5, leading=14, textColor=ink)
    h1 = ParagraphStyle("H1", parent=styles["Title"], fontSize=26, leading=32, textColor=ink, spaceAfter=4)
    sub = ParagraphStyle("Sub", parent=body, fontSize=13, textColor=muted)
    brand = ParagraphStyle("Brand", parent=body, fontSize=9, textColor=accent, spaceAfter=10)
    sec = ParagraphStyle("Sec", parent=styles["Heading2"], fontSize=14, leading=18, textColor=ink, spaceBefore=14, spaceAfter=8)
    callout = ParagraphStyle("Callout", parent=body, leftIndent=10, textColor=muted, fontName="Helvetica-Oblique")
    toc_item = ParagraphStyle("TocItem", parent=body, fontSize=10.5, leading=17)

    def _header_footer(canvas, docT):  # noqa: ANN001
        canvas.saveState()
        w, h = A4
        canvas.setStrokeColor(line)
        canvas.line(16 * mm, 12 * mm, w - 16 * mm, 12 * mm)
        canvas.setFont("Helvetica", 7.5)
        canvas.setFillColor(muted)
        canvas.drawString(16 * mm, 8 * mm, f"{doc['brand']} · {doc['title']} · {doc['subtitle']}")
        canvas.drawRightString(w - 16 * mm, 8 * mm, f"Page {canvas.getPageNumber()}")
        canvas.restoreState()

    buf = io.BytesIO()
    pdf = SimpleDocTemplate(
        buf, pagesize=A4, leftMargin=16 * mm, rightMargin=16 * mm, topMargin=18 * mm, bottomMargin=18 * mm,
        title=f"{doc['title']} — {doc['subtitle']}", author=doc["brand"],
    )

    story: list[Any] = [
        Paragraph(doc["brand"].upper(), brand),
        Paragraph(doc["title"], h1),
        Paragraph(doc["subtitle"], sub),
        Spacer(1, 6),
        Paragraph(f"Generated {_date(doc)} · Confidential — internal use", ParagraphStyle("m", parent=body, textColor=muted, fontSize=8.5)),
        Spacer(1, 18),
        Paragraph("Table of Contents", sec),
    ]
    for i, s in enumerate(doc["sections"], start=1):
        story.append(Paragraph(f"{i:02d}&nbsp;&nbsp;{s['heading']}", toc_item))
    story.append(PageBreak())

    table_style = TableStyle(
        [
            ("BACKGROUND", (0, 0), (-1, 0), bg),
            ("TEXTCOLOR", (0, 0), (-1, 0), ink),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("GRID", (0, 0), (-1, -1), 0.5, line),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#FBFCFE")]),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]
    )
    kv_style = TableStyle(
        [
            ("BACKGROUND", (0, 0), (0, -1), bg),
            ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 8.5),
            ("GRID", (0, 0), (-1, -1), 0.5, line),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]
    )

    usable = A4[0] - 32 * mm
    for i, s in enumerate(doc["sections"], start=1):
        story.append(Paragraph(f"<font color='{ACCENT}'>{i:02d}</font>&nbsp;&nbsp;{s['heading']}", sec))
        for block in s["blocks"]:
            kind = block.get("type")
            if kind == "p" and block.get("text"):
                story.append(Paragraph(_esc(block["text"]), body))
            elif kind == "callout" and block.get("text"):
                quote = Table([[Paragraph(_esc(block["text"]), callout)]], colWidths=[usable])
                quote.setStyle(
                    TableStyle(
                        [
                            ("BACKGROUND", (0, 0), (-1, -1), bg),
                            ("LINEBEFORE", (0, 0), (0, -1), 2.5, accent),
                            ("TOPPADDING", (0, 0), (-1, -1), 8),
                            ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
                            ("LEFTPADDING", (0, 0), (-1, -1), 10),
                        ]
                    )
                )
                story.append(quote)
                story.append(Spacer(1, 6))
            elif kind == "list":
                for j, item in enumerate(block.get("items", []), start=1):
                    bullet = f"{j}." if block.get("ordered") else "•"
                    story.append(Paragraph(f"{bullet}&nbsp;&nbsp;{_esc(str(item))}", body))
                story.append(Spacer(1, 6))
            elif kind == "kv" and block.get("items"):
                rows = [[Paragraph(_esc(str(k)), body), Paragraph(_esc(str(v)), body)] for k, v in block["items"]]
                t = Table(rows, colWidths=[usable * 0.34, usable * 0.66])
                t.setStyle(kv_style)
                story.append(t)
                story.append(Spacer(1, 6))
            elif kind == "table" and block.get("rows"):
                cols = block.get("columns", [])
                data = [[Paragraph(f"<b>{_esc(str(c))}</b>", body) for c in cols]]
                data += [[Paragraph(_esc(str(c)), body) for c in row] for row in block["rows"][:30]]
                t = Table(data, colWidths=[usable / max(1, len(cols))] * len(cols), repeatRows=1)
                t.setStyle(table_style)
                story.append(t)
                story.append(Spacer(1, 6))
        story.append(Spacer(1, 8))

    pdf.build(story, onFirstPage=_header_footer, onLaterPages=_header_footer)
    return buf.getvalue()


def _esc(text: str) -> str:
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


# --- Word (python-docx) --------------------------------------------------------
def render_docx(doc: dict[str, Any]) -> bytes:
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Pt, RGBColor
    except ImportError as exc:  # pragma: no cover
        raise ExportUnavailable("Word export requires the 'python-docx' package.") from exc

    accent = RGBColor(0x25, 0x63, 0xEB)
    muted = RGBColor(0x5B, 0x6B, 0x84)

    d = Document()
    styles = d.styles
    styles["Normal"].font.name = "Calibri"
    styles["Normal"].font.size = Pt(10.5)

    # Cover
    p = d.add_paragraph()
    run = p.add_run(doc["brand"].upper())
    run.font.color.rgb = accent
    run.font.bold = True
    run.font.size = Pt(9)
    title = d.add_paragraph()
    run = title.add_run(doc["title"])
    run.font.size = Pt(30)
    run.font.bold = True
    sub = d.add_paragraph()
    run = sub.add_run(doc["subtitle"])
    run.font.size = Pt(14)
    run.font.color.rgb = muted
    meta = d.add_paragraph()
    run = meta.add_run(f"Generated {_date(doc)} · Confidential — internal use")
    run.font.size = Pt(8.5)
    run.font.color.rgb = muted

    # Footer with page numbers
    section = d.sections[0]
    footer_p = section.footer.paragraphs[0]
    footer_p.text = f"{doc['brand']} · {doc['title']}\t\t"
    footer_p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    _docx_page_number(footer_p)
    for r in footer_p.runs:
        r.font.size = Pt(8)
        r.font.color.rgb = muted

    # TOC
    d.add_heading("Table of Contents", level=1)
    for i, s in enumerate(doc["sections"], start=1):
        d.add_paragraph(f"{i:02d}   {s['heading']}")
    d.add_page_break()

    for i, s in enumerate(doc["sections"], start=1):
        d.add_heading(f"{i:02d}  {s['heading']}", level=1)
        for block in s["blocks"]:
            kind = block.get("type")
            if kind == "p" and block.get("text"):
                d.add_paragraph(block["text"])
            elif kind == "callout" and block.get("text"):
                quote = d.add_paragraph(block["text"], style="Intense Quote")
                for r in quote.runs:
                    r.font.color.rgb = muted
            elif kind == "list":
                style = "List Number" if block.get("ordered") else "List Bullet"
                for item in block.get("items", []):
                    d.add_paragraph(str(item), style=style)
            elif kind == "kv" and block.get("items"):
                t = d.add_table(rows=0, cols=2)
                t.style = "Light Grid Accent 1"
                for k, v in block["items"]:
                    row = t.add_row().cells
                    row[0].text = str(k)
                    row[0].paragraphs[0].runs[0].font.bold = True
                    row[1].text = str(v)
                d.add_paragraph()
            elif kind == "table" and block.get("rows"):
                cols = block.get("columns", [])
                t = d.add_table(rows=1, cols=len(cols))
                t.style = "Light Grid Accent 1"
                for j, c in enumerate(cols):
                    cell = t.rows[0].cells[j]
                    cell.text = str(c)
                    cell.paragraphs[0].runs[0].font.bold = True
                for row in block["rows"][:30]:
                    cells = t.add_row().cells
                    for j, val in enumerate(row[: len(cols)]):
                        cells[j].text = str(val)
                d.add_paragraph()

    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def _docx_page_number(paragraph) -> None:  # noqa: ANN001
    """Append a live PAGE field to a paragraph (python-docx has no helper)."""
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    run = paragraph.add_run()
    for tag, attrs, text in (
        ("w:fldChar", {"w:fldCharType": "begin"}, None),
        ("w:instrText", {"xml:space": "preserve"}, " PAGE "),
        ("w:fldChar", {"w:fldCharType": "end"}, None),
    ):
        el = OxmlElement(tag)
        for k, v in attrs.items():
            el.set(qn(k), v)
        if text:
            el.text = text
        run._r.append(el)  # noqa: SLF001


# --- PowerPoint (python-pptx) ---------------------------------------------------
def render_pptx(doc: dict[str, Any]) -> bytes:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except ImportError as exc:  # pragma: no cover
        raise ExportUnavailable("PowerPoint export requires the 'python-pptx' package.") from exc

    accent = RGBColor(0x25, 0x63, 0xEB)
    ink = RGBColor(0x16, 0x23, 0x3B)
    muted = RGBColor(0x5B, 0x6B, 0x84)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_slide():
        return prs.slides.add_slide(blank)

    def textbox(slide, left, top, width, height):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        box.text_frame.word_wrap = True
        return box.text_frame

    # Cover slide
    slide = add_slide()
    tf = textbox(slide, 0.8, 2.2, 11.7, 3)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = doc["brand"].upper()
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = accent
    p2 = tf.add_paragraph()
    run = p2.add_run()
    run.text = doc["title"]
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = ink
    p3 = tf.add_paragraph()
    run = p3.add_run()
    run.text = f"{doc['subtitle']} · Generated {_date(doc)}"
    run.font.size = Pt(16)
    run.font.color.rgb = muted

    # Agenda slide
    slide = add_slide()
    tf = textbox(slide, 0.8, 0.6, 11.7, 0.9)
    run = tf.paragraphs[0].add_run()
    run.text = "Agenda"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = ink
    tf = textbox(slide, 0.8, 1.6, 11.7, 5.3)
    for i, s in enumerate(doc["sections"], start=1):
        p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"{i:02d}   {s['heading']}"
        run.font.size = Pt(16)
        run.font.color.rgb = ink

    # One slide per section
    for i, s in enumerate(doc["sections"], start=1):
        slide = add_slide()
        tf = textbox(slide, 0.8, 0.5, 11.7, 0.9)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{i:02d}  "
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = accent
        run = p.add_run()
        run.text = s["heading"]
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = ink

        table_block = next((b for b in s["blocks"] if b.get("type") in ("table", "kv")), None)
        text_blocks = [b for b in s["blocks"] if b.get("type") in ("p", "callout", "list")]

        body_width = 5.9 if table_block else 11.7
        tf = textbox(slide, 0.8, 1.5, body_width, 5.4)
        first = True
        for block in text_blocks[:3]:
            if block["type"] in ("p", "callout") and block.get("text"):
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                run = p.add_run()
                run.text = block["text"][:400]
                run.font.size = Pt(13)
                run.font.color.rgb = muted if block["type"] == "callout" else ink
                run.font.italic = block["type"] == "callout"
            elif block["type"] == "list":
                for item in block.get("items", [])[:7]:
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    run = p.add_run()
                    run.text = f"• {str(item)[:160]}"
                    run.font.size = Pt(13)
                    run.font.color.rgb = ink

        if table_block:
            if table_block["type"] == "kv":
                cols = ["", ""]
                rows = table_block.get("items", [])[:8]
            else:
                cols = table_block.get("columns", [])[:5]
                rows = [r[: len(cols)] for r in table_block.get("rows", [])[:8]]
            if rows:
                shape = slide.shapes.add_table(
                    len(rows) + (1 if table_block["type"] == "table" else 0),
                    max(2, len(cols)),
                    Inches(7.0),
                    Inches(1.5),
                    Inches(5.5),
                    Inches(0.4 * (len(rows) + 1)),
                )
                table = shape.table
                offset = 0
                if table_block["type"] == "table":
                    for j, c in enumerate(cols):
                        cell = table.cell(0, j)
                        cell.text = str(c)
                        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
                        cell.text_frame.paragraphs[0].runs[0].font.bold = True
                    offset = 1
                for r, row in enumerate(rows):
                    for j, val in enumerate(row):
                        cell = table.cell(r + offset, j)
                        cell.text = str(val)[:80]
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.size = Pt(10.5)

        # Slide footer
        tf = textbox(slide, 0.8, 7.05, 11.7, 0.35)
        run = tf.paragraphs[0].add_run()
        run.text = f"{doc['brand']} · {doc['title']} · {i + 2}"
        run.font.size = Pt(9)
        run.font.color.rgb = muted

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
